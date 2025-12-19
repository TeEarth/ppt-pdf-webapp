import os
import pandas as pd
from pptx import Presentation
from pdf2image import convert_from_path

def POST(request):
    form = request.form()
    names = [form.get(f"name{i}", "") for i in range(1, 6)]

    sheet_url = "https://docs.google.com/spreadsheets/d/17XiGo7Q4VFJYAZ4ZydIwLOC1SxUXthZB5J8Jb_mYX-U/export?format=csv"
    df = pd.read_csv(sheet_url)

    prs = Presentation("public/template/template.pptx")

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for i, name in enumerate(names, start=1):
                    shape.text = shape.text.replace(f"{{{{Name{i}}}}}", name)

    ppt_path = "/tmp/output.pptx"
    prs.save(ppt_path)

    images = convert_from_path(ppt_path)
    output = []
    for i, img in enumerate(images, 1):
        path = f"/tmp/slide{i}.png"
        img.save(path)
        output.append(path)

    return {
        "ppt": ppt_path,
        "images": output
    }
